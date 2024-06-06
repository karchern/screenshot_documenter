from documenter.coordinates import BoxCoordinates
from documenter.listeners import Listeners
from documenter.grab_screenshot import get_screenshot_from_screen
import logging
import tkinter as tk
from tkinter import simpledialog
from tkinter.filedialog import asksaveasfile

from pptx import Presentation 
from pptx.util import Inches
import io

class Runner:
    def __init__(self):
        self.screenshot_listeners = None
        self.breakout_listener = Listeners()
        self.breakout_listener.add_listener_and_start('breakout_listener', self.breakout_listener.instantiate_breakout_listener)
        self.screenshots = []

    def get_output_file(self):
        root = tk.Tk()
        root.withdraw()
        file = asksaveasfile(mode='w', defaultextension=".txt")
        return file

    def save_screenshots_as_powerpoint(self, output_file_path):

        if not output_file_path.endswith('.pptx'):
            output_file_path = output_file_path.split('.')[:-1]
            output_file_path += '.pptx'
        prs = Presentation() 
        blank_slide_layout = prs.slide_layouts[6] 
        left = top = Inches(0)        
        for screenshot in self.screenshots:
            slide_height = 7.5
            slide_width = 10            
            aspect_ratio = screenshot.size[0] / screenshot.size[1]
            if aspect_ratio > 1.25:
                slide_height = 10 / aspect_ratio
            else:
                slide_width = 7.5 * aspect_ratio

            slide = prs.slides.add_slide(blank_slide_layout)
            with io.BytesIO() as output:
                screenshot.save(output, format='PNG')       
                pic = slide.shapes.add_picture(output, left, top, height=Inches(slide_height), width=Inches(slide_width))
        prs.save(output_file_path)
        return 

    def run(self):
        while True:
        # Instantiate listener if not already/
            if not self.screenshot_listeners:
                self.screenshot_listeners = Listeners()
                self.screenshot_listeners.add_listener_and_start('mouse_listener_terminator', self.screenshot_listeners.instantiate_terminator_mouse_listener)
                self.screenshot_listeners.add_listener_and_start('keyboard_listener', self.screenshot_listeners.instantiate_keyboard_listener_screenshot)

            
            if isinstance(self.screenshot_listeners.mouse_coordinates, BoxCoordinates) and self.screenshot_listeners:
                logging.info('Rectangle coordinates received, taking screenshot...')
                screenshot = get_screenshot_from_screen(self.screenshot_listeners.mouse_coordinates)
                self.screenshots.append(screenshot)
                self.screenshot_listeners = None
            if not self.breakout_listener.breakout_listener.running:
                break
        output_file = self.get_output_file()
        self.save_screenshots_as_powerpoint(output_file.name)

            




    

